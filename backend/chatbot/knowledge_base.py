import os
from typing import List, Dict
import logging
import requests
from bs4 import BeautifulSoup
import PyPDF2
import io
from pptx import Presentation
import json
import hashlib

logger = logging.getLogger(__name__)

# Simple in-memory vector store for deployment compatibility
# Uses OpenAI embeddings via Emergent LLM key instead of sentence-transformers

class SimpleVectorStore:
    """Lightweight vector store using keyword matching and TF-IDF-like scoring"""
    
    def __init__(self):
        self.documents = {}  # id -> {text, metadata, keywords}
        self.storage_path = "/app/backend/knowledge_store.json"
        self._load_store()
    
    def _load_store(self):
        """Load documents from disk"""
        try:
            if os.path.exists(self.storage_path):
                with open(self.storage_path, 'r') as f:
                    self.documents = json.load(f)
                logger.info(f"Loaded {len(self.documents)} documents from store")
        except Exception as e:
            logger.warning(f"Could not load store: {e}")
            self.documents = {}
    
    def _save_store(self):
        """Save documents to disk"""
        try:
            with open(self.storage_path, 'w') as f:
                json.dump(self.documents, f)
        except Exception as e:
            logger.error(f"Could not save store: {e}")
    
    def _extract_keywords(self, text: str) -> List[str]:
        """Extract important keywords from text"""
        # Remove common stop words
        stop_words = {
            'the', 'a', 'an', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for',
            'of', 'with', 'by', 'from', 'as', 'is', 'was', 'are', 'were', 'been',
            'be', 'have', 'has', 'had', 'do', 'does', 'did', 'will', 'would',
            'could', 'should', 'may', 'might', 'must', 'shall', 'can', 'need',
            'it', 'its', 'this', 'that', 'these', 'those', 'i', 'you', 'he',
            'she', 'we', 'they', 'what', 'which', 'who', 'when', 'where', 'why',
            'how', 'all', 'each', 'every', 'both', 'few', 'more', 'most', 'other',
            'some', 'such', 'no', 'nor', 'not', 'only', 'own', 'same', 'so',
            'than', 'too', 'very', 'just', 'our', 'your', 'their', 'my', 'his', 'her'
        }
        
        words = text.lower().split()
        # Clean words and filter
        keywords = []
        for word in words:
            clean = ''.join(c for c in word if c.isalnum())
            if clean and len(clean) > 2 and clean not in stop_words:
                keywords.append(clean)
        
        return list(set(keywords))
    
    def add(self, doc_id: str, text: str, metadata: Dict = None):
        """Add a document to the store"""
        keywords = self._extract_keywords(text)
        self.documents[doc_id] = {
            'text': text,
            'metadata': metadata or {},
            'keywords': keywords
        }
        self._save_store()
    
    def query(self, question: str, n_results: int = 3) -> List[str]:
        """Query documents using keyword matching"""
        if not self.documents:
            return []
        
        query_keywords = set(self._extract_keywords(question))
        if not query_keywords:
            return []
        
        # Score each document by keyword overlap
        scores = []
        for doc_id, doc in self.documents.items():
            doc_keywords = set(doc['keywords'])
            overlap = len(query_keywords & doc_keywords)
            if overlap > 0:
                # Score = overlap / query length (for relevance)
                score = overlap / len(query_keywords)
                scores.append((score, doc['text']))
        
        # Sort by score and return top results
        scores.sort(reverse=True, key=lambda x: x[0])
        return [text for _, text in scores[:n_results]]
    
    def clear(self):
        """Clear all documents"""
        self.documents = {}
        self._save_store()


class KnowledgeBase:
    def __init__(self):
        # Use simple vector store for deployment compatibility
        self.store = SimpleVectorStore()
        logger.info("Knowledge Base initialized (deployment-ready mode)")
    
    def add_document(self, doc_id: str, text: str, metadata: Dict = None):
        """Add a document to the knowledge base"""
        try:
            # Split text into chunks for better retrieval
            chunks = self._chunk_text(text)
            
            for i, chunk in enumerate(chunks):
                chunk_id = f"{doc_id}_chunk_{i}"
                self.store.add(chunk_id, chunk, metadata)
            
            logger.info(f"Added document {doc_id} with {len(chunks)} chunks")
            return True
        except Exception as e:
            logger.error(f"Error adding document: {e}")
            return False
    
    def query(self, question: str, n_results: int = 3) -> List[str]:
        """Query the knowledge base"""
        try:
            results = self.store.query(question, n_results)
            return results
        except Exception as e:
            logger.error(f"Error querying knowledge base: {e}")
            return []
    
    def _chunk_text(self, text: str, chunk_size: int = 500) -> List[str]:
        """Split text into chunks"""
        words = text.split()
        chunks = []
        current_chunk = []
        current_size = 0
        
        for word in words:
            current_chunk.append(word)
            current_size += len(word) + 1
            
            if current_size >= chunk_size:
                chunks.append(' '.join(current_chunk))
                current_chunk = []
                current_size = 0
        
        if current_chunk:
            chunks.append(' '.join(current_chunk))
        
        return chunks
    
    def crawl_website(self, base_url: str) -> bool:
        """Crawl website pages and add to knowledge base"""
        try:
            pages = [
                f"{base_url}/",
                f"{base_url}/about",
                f"{base_url}/services",
                f"{base_url}/products",
            ]
            
            for page_url in pages:
                try:
                    response = requests.get(page_url, timeout=10)
                    if response.status_code == 200:
                        soup = BeautifulSoup(response.content, 'html.parser')
                        
                        # Remove script and style elements
                        for script in soup(["script", "style"]):
                            script.decompose()
                        
                        # Get text
                        text = soup.get_text()
                        lines = (line.strip() for line in text.splitlines())
                        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
                        text = '\n'.join(chunk for chunk in chunks if chunk)
                        
                        # Create doc_id from URL
                        doc_id = hashlib.md5(page_url.encode()).hexdigest()[:12]
                        
                        self.add_document(
                            doc_id=doc_id,
                            text=text,
                            metadata={"source": page_url, "type": "webpage"}
                        )
                        
                        logger.info(f"Crawled and added: {page_url}")
                except Exception as e:
                    logger.error(f"Error crawling {page_url}: {e}")
                    continue
            
            return True
        except Exception as e:
            logger.error(f"Error in website crawl: {e}")
            return False
    
    def process_pdf(self, file_content: bytes, filename: str) -> bool:
        """Process a PDF file and add to knowledge base"""
        try:
            pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
            text = ""
            
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
            
            doc_id = hashlib.md5(filename.encode()).hexdigest()[:12]
            return self.add_document(
                doc_id=doc_id,
                text=text,
                metadata={"source": filename, "type": "pdf"}
            )
        except Exception as e:
            logger.error(f"Error processing PDF: {e}")
            return False
    
    def process_pptx(self, file_content: bytes, filename: str) -> bool:
        """Process a PowerPoint file and add to knowledge base"""
        try:
            prs = Presentation(io.BytesIO(file_content))
            text = ""
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            
            doc_id = hashlib.md5(filename.encode()).hexdigest()[:12]
            return self.add_document(
                doc_id=doc_id,
                text=text,
                metadata={"source": filename, "type": "pptx"}
            )
        except Exception as e:
            logger.error(f"Error processing PPTX: {e}")
            return False
    
    def get_stats(self) -> Dict:
        """Get knowledge base statistics"""
        return {
            "total_documents": len(self.store.documents),
            "storage_mode": "keyword-based (deployment-ready)"
        }

# Global instance
knowledge_base = KnowledgeBase()
